#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Criar apresentação
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Definir cores
DARK_BLUE = RGBColor(15, 23, 42)
ACCENT_RED = RGBColor(220, 38, 38)
LIGHT_GRAY = RGBColor(248, 250, 252)
TEXT_GRAY = RGBColor(75, 85, 99)

def add_title_slide(prs, title, subtitle):
    """Adiciona um slide de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout em branco
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_RED
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_list):
    """Adiciona um slide com conteúdo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY

    # Barra superior
    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = DARK_BLUE
    top_bar.line.color.rgb = ACCENT_RED
    top_bar.line.width = Pt(3)

    # Título
    title_frame = top_bar.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(8)

    # Conteúdo
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(6))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, item in enumerate(content_list):
        if i > 0:
            content_frame.add_paragraph()

        p = content_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_GRAY
        p.space_after = Pt(16)
        p.level = 0

    return slide

# Slide 1: Capa
add_title_slide(prs, "QUIZ TRÂNSITO", "Sistema de Aprendizado para Habilitação")

# Slide 2: Alunos
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = LIGHT_GRAY

top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = DARK_BLUE
top_bar.line.color.rgb = ACCENT_RED
top_bar.line.width = Pt(3)

title_frame = top_bar.text_frame
p = title_frame.paragraphs[0]
p.text = "Integrantes do Projeto"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.space_before = Pt(8)

alunos = [
    "Roberto Gustavo Reiniger Neto - RA: 24209569",
    "Paulo Heredia Padin Junior - RA: 24202269",
    "Humberto Assis Andrade Paixão - RA: 2229630",
    "Belchior Rodrigues Cavalcante - RA: 2231611",
    "Shirley Christina Ramos dos Santos - RA: 2213441",
    "Guilherme Campbell Penna - RA: 24216190",
    "Arthur Ferreira de Castro - RA: 24228235"
]

alunos_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(5.5))
alunos_frame = alunos_box.text_frame
alunos_frame.word_wrap = True

for i, aluno in enumerate(alunos):
    if i > 0:
        alunos_frame.add_paragraph()
    p = alunos_frame.paragraphs[i]
    p.text = "• " + aluno
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_GRAY
    p.space_after = Pt(10)

# Slide 3: Escopo do Projeto
add_content_slide(prs, "Escopo do Projeto", [
    "Framework Web: React (Frontend) + Node.js/Express (Backend)",
    "Controle de Versao: Git/GitHub",
    "Banco de Dados: MySQL (SQL)",
    "Funcionalidades: Quiz interativo, informacoes sobre habilitacao, painel administrativo"
])

# Slide 4: Tecnologias
add_content_slide(prs, "Tecnologias Utilizadas", [
    "Frontend: React 18 + Bootstrap 5 + Axios",
    "Backend: Node.js + Express.js",
    "Banco de Dados: MySQL com Promises",
    "Autenticação: JWT (JSON Web Tokens)",
    "Upload de Imagens: Multer"
])

# Slide 5: Funcionalidades Principais
add_content_slide(prs, "Funcionalidades Principais", [
    "Quiz Interativo: Simulados com seleção aleatória de perguntas",
    "Gestão de Perguntas: CRUD completo com upload de imagens",
    "Material de Estudo: Informações estruturadas sobre habilitação",
    "Painel Administrativo: Controle de conteúdo e permissões",
    "Histórico de Resultados: Armazenamento de desempenho"
])

# Slide 6: Novas Regras de Trânsito 2024
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = LIGHT_GRAY

top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = DARK_BLUE
top_bar.line.color.rgb = ACCENT_RED
top_bar.line.width = Pt(3)

title_frame = top_bar.text_frame
p = title_frame.paragraphs[0]
p.text = "Novas Regras de Trânsito - 2024"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.space_before = Pt(8)

regras = [
    "Validade da CNH: Aumentada para 10 anos (menores de 50 anos)",
    "Pontuação: 40 pontos para condutores sem infrações gravíssimas",
    "Curso de Reciclagem: Obrigatório ao atingir limite de pontos",
    "Exame Toxicológico: Obrigatório para categorias C, D e E",
    "CNH Digital: Versão digital com valor legal igual à física"
]

regras_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
regras_frame = regras_box.text_frame
regras_frame.word_wrap = True

for i, regra in enumerate(regras):
    if i > 0:
        regras_frame.add_paragraph()
    p = regras_frame.paragraphs[i]
    p.text = "• " + regra
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_GRAY
    p.space_after = Pt(12)

# Slide 7: Mudanças Importantes nas Infrações
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = LIGHT_GRAY

top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = DARK_BLUE
top_bar.line.color.rgb = ACCENT_RED
top_bar.line.width = Pt(3)

title_frame = top_bar.text_frame
p = title_frame.paragraphs[0]
p.text = "Mudanças em Infrações e Penalidades"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.space_before = Pt(8)

mudancas = [
    "Dirigir Sem Habilitação: Multa aumentada para R$ 11.769,50",
    "Uso de Celular: Agora faz perder 5 pontos (antes era 4)",
    "Estacionamento em Vaga PcD: Multa de R$ 5.884,50",
    "Buzina Excessiva: Infração gravíssima em áreas sensíveis",
    "Excesso de Velocidade: Multa aumentada em até 50% dependendo do excesso"
]

mudancas_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
mudancas_frame = mudancas_box.text_frame
mudancas_frame.word_wrap = True

for i, mudanca in enumerate(mudancas):
    if i > 0:
        mudancas_frame.add_paragraph()
    p = mudancas_frame.paragraphs[i]
    p.text = "• " + mudanca
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_GRAY
    p.space_after = Pt(12)

# Slide 8: Encerramento
add_title_slide(prs, "Demonstração", "Aplicação em Funcionamento")

# Salvar apresentação
prs.save(r'c:\xampp\htdocs\quiz-transito\apresentacao_quiz_transito.pptx')
print("Apresentacao criada com sucesso: apresentacao_quiz_transito.pptx")
